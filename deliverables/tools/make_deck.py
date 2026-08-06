#!/usr/bin/env python3
"""Build deliverables/Supplementary_Deck_EN.pptx (16:9, English)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "deliverables" / "Supplementary_Deck_EN.pptx"
ASSETS = ROOT / "deliverables" / "assets"
BAR = ROOT / "ComfyUI" / "output" / "废土酒吧环境"

BG = RGBColor(0x14, 0x16, 0x10)
PANEL = RGBColor(0x1E, 0x22, 0x18)
ACCENT = RGBColor(0xE0, 0xA6, 0x3C)
TEXT = RGBColor(0xE8, 0xE4, 0xD8)
DIM = RGBColor(0x9A, 0x94, 0x84)
GOOD = RGBColor(0x7D, 0xBD, 0x5A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(1, l, t, w, h)  # rectangle
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(s, l, t, w, h, runs, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, size, color, bold)]
    first = True
    for rtext, rsize, rcolor, rbold in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = rtext
        r.font.size = Pt(rsize); r.font.color.rgb = rcolor; r.font.bold = rbold
        r.font.name = "Calibri"
    return tb


def bullets(s, l, t, w, h, items, size=15, gap=8, color=TEXT, bold_head=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▪  " + it
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
        if bold_head: r.font.bold = True
    return tb


def kicker(s, text):
    txt(s, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4), text.upper(), size=13, color=ACCENT, bold=True)


def header(s, title, sub=None):
    kicker(s, "IndieGen Asset Studio · Track 1")
    txt(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.9), title, size=30, color=TEXT, bold=True)
    if sub:
        txt(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.5), sub, size=15, color=DIM)
    box(s, Inches(0.7), Inches(2.05), Inches(11.9), Emu(9525), fill=ACCENT)


def pic(s, path, l, t, w=None, h=None):
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    return s.shapes.add_picture(str(path), l, t, **kw)


# ---------------------------------------------------------------- 1 title
s = slide()
box(s, 0, 0, prs.slide_width, Inches(0.18), fill=ACCENT)
txt(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.2), "IndieGen Asset Studio", size=48, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.7),
    "Track 1 · Multimodal Content Creation Tools", size=24, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(3.9), Inches(11.3), Inches(0.9),
    "From one text setting to a complete, style-consistent asset chain: concept art, depth, PBR materials, "
    "music, ambient SFX, sprite atlas and a 3D character — natively on AMD Radeon / ROCm.",
    size=16, color=DIM, align=PP_ALIGN.CENTER)
box(s, Inches(4.35), Inches(5.6), Inches(4.6), Inches(0.65), fill=PANEL, line=ACCENT)
txt(s, Inches(4.35), Inches(5.7), Inches(4.6), Inches(0.5),
    "AMD Radeon gfx1100 · ROCm 7.1 · ComfyUI 0.30", size=14, color=GOOD, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- 2 problem & solution
s = slide(); header(s, "Problem & Solution", "Indie teams spend 40-70% of budget on outsourced art — and style drifts between vendors.")
bullets(s, Inches(0.7), Inches(2.5), Inches(6.2), Inches(4.2), [
    "One chain instead of many vendors — every asset is conditioned on the same character sheet / setting image.",
    "One GPU instead of a cloud bill — the whole chain runs locally on a Radeon card through ROCm.",
    "Engine-ready outputs — PNG atlas + JSON, PBR maps, GLB/OBJ, MP3 loops, no post-processing.",
    "Human-in-the-loop + full reproducibility — every job stores prompt, seed, steps, CFG and model.",
    "Bilingual UX — Chinese prompts are translated offline (MarianMT) for SDXL / MusicGen.",
], size=16)
pic(s, ASSETS / "pipeline.png", Inches(7.0), Inches(2.5), w=Inches(5.7))

# ---------------------------------------------------------------- 3 pipeline
s = slide(); header(s, "The Full Asset Pipeline", "W1 → W7: one text setting drives seven asset types.")
pic(s, ASSETS / "pipeline.png", Inches(1.2), Inches(2.4), w=Inches(10.9))

# ---------------------------------------------------------------- 4 architecture
s = slide(); header(s, "System Architecture", "Browser UI · FastAPI backend · ComfyUI GPU engine.")
pic(s, ASSETS / "architecture.png", Inches(1.5), Inches(2.4), w=Inches(10.3))

# ---------------------------------------------------------------- 5 models & algorithms
s = slide(); header(s, "Models & Algorithms")
rows = [
    ("Stage", "Method / model", "Algorithm notes"),
    ("W1 concept", "SDXL base 1.0", "Text → latent → KSampler (Euler, fixed seed) → VAE decode"),
    ("W2 depth", "MiDaS DPT-Hybrid", "Monocular depth via dense prediction transformer"),
    ("W3 PBR", "Custom PBRFromDepth", "Depth gradients → normal/height/roughness/metalness"),
    ("W4 music", "MusicGen-small", "Autoregressive audio conditioned on text prompt, 44.1 kHz"),
    ("W5 SFX", "Custom ProceduralSFX", "Glass clinks, murmur, bar ambience (DSP synthesis)"),
    ("W6 sprite", "ActionPose + ControlNet OpenPose (SDXL)", "18-keypoint skeletons → per-frame img2img → atlas + JSON"),
    ("W7 2D→3D", "TripoSR", "Single image → triplane → 3D mesh (GLB/OBJ)"),
    ("Translate", "MarianMT opus-mt-zh-en", "Offline Chinese→English prompt translation"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(2.35), Inches(11.9), Inches(4.5))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.7); tbl.columns[1].width = Inches(4.0); tbl.columns[2].width = Inches(6.2)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if r else RGBColor(0x2A, 0x2E, 0x22)
        cell.margin_top = cell.margin_bottom = Pt(4)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; rn = p.add_run(); rn.text = val
        rn.font.size = Pt(12); rn.font.name = "Calibri"
        rn.font.color.rgb = ACCENT if (r == 0 or c == 0) else TEXT
        rn.font.bold = (r == 0 or c == 0)

# ---------------------------------------------------------------- 6 AMD ROCm
s = slide(); header(s, "AMD Radeon GPU / ROCm Adaptation", "Auto-detected ROCm wheel install · measured live on gfx1100, 48 GiB VRAM.")
bullets(s, Inches(0.7), Inches(2.45), Inches(6.0), Inches(3.6), [
    "setup_comfyui.sh auto-detects GPU: NVIDIA → CUDA 12.8, AMD (/dev/kfd) → ROCm 7.1 wheels, CPU fallback — matching torch/torchvision/torchaudio pins.",
    "PyTorch's ROCm build exposes the standard cuda API, so ComfyUI, transformers, ControlNet and all custom nodes run unchanged.",
    "rocm-smi sampled every second during the benchmark: GPU use, VRAM, temperature and power.",
], size=15)
pic(s, ASSETS / "performance.png", Inches(6.9), Inches(2.45), w=Inches(5.9))
box(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.8), fill=PANEL)
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5),
    "Full chain W1–W6: ≈57 s GPU execution · W7 TripoSR 2D→3D: 5.2 s wall on the live run",
    size=15, color=GOOD, bold=True)

# ---------------------------------------------------------------- 7 demo scenario
s = slide(); header(s, "Creative Scenario — The Wasteland Bartender")
concept = BAR / "W1_concept_00001_.png"
depth = BAR / "W2_depth_00001_.png"
atlas = BAR / "W6_atlas_00002_.png"
pbr = BAR / "W3_normal_00001_.png"
pic(s, concept, Inches(0.7), Inches(2.3), h=Inches(3.4))
pic(s, depth, Inches(4.0), Inches(2.3), h=Inches(3.4))
pic(s, pbr, Inches(7.3), Inches(2.3), h=Inches(3.4))
if atlas.exists():
    pic(s, atlas, Inches(10.3), Inches(2.3), h=Inches(3.4))
txt(s, Inches(0.7), Inches(5.9), Inches(12), Inches(1.2),
    "One text setting → concept art (15 s) · depth + PBR maps · dark industrial music · glass-clink SFX · "
    "8-frame run atlas + JSON (21 s) · TripoSR 3D character (5.2 s) — assembled into an auto-animated 3D bar.",
    size=14, color=DIM)

# ---------------------------------------------------------------- 8 value
s = slide(); header(s, "Innovation & Practical Value")
bullets(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(4.4), [
    "Full-chain style lock: one conditioning anchor across seven asset types — no cross-vendor drift.",
    "Engine-ready outputs: PNG atlas + JSON frame config, PBR maps, GLB/OBJ — import directly.",
    "Low-cost local GPU: peak VRAM 21.7 GB on gfx1100; weight offloading supports 8–16 GB Radeon cards.",
    "Reproducible & auditable: every job stores its full parameter set; fixed seeds repeat results.",
    "Product-grade packaging: docker-compose + nginx + systemd, proxy-only networking (8188 never exposed).",
    "A living demo: the 3D scene assembles a batch into an automated bartender performance with generated music/SFX.",
], size=16)

# ---------------------------------------------------------------- 9 close
s = slide(); header(s, "Deliverables & How to Run")
bullets(s, Inches(0.7), Inches(2.5), Inches(7.0), Inches(4.0), [
    "Project profile (PDF) · this deck · 3–5 min demo video · live benchmark data.",
    "Full source: workflows W1–W7, 5 custom ComfyUI nodes, FastAPI backend + SQLite, browser UI + Three.js scene.",
    "Run: bash scripts/setup_comfyui.sh && bash scripts/download_models.sh && bash server/run.sh",
    "W7 3D: bash scripts/setup_3d.sh && bash scripts/download_3d_model.sh",
    "Benchmark: ./ComfyUI/.venv/bin/python scripts/benchmark_amd.py",
], size=16)
box(s, Inches(8.0), Inches(2.5), Inches(4.6), Inches(3.4), fill=PANEL, line=ACCENT)
txt(s, Inches(8.3), Inches(2.8), Inches(4.0), Inches(3.0), [
    ("Track 1 · Multimodal Content Creation Tools", 17, ACCENT, True),
    ("\nOne text setting → game-ready multimodal assets, locally on AMD Radeon / ROCm.", 14, DIM, False),
    ("\n\ngithub.com/MFWTW/create_3Dgame", 14, GOOD, True),
], spacing=1.15)

prs.save(str(OUT))
print("PPTX written:", OUT, OUT.stat().st_size, "bytes")
